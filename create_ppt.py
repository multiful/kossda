"""
KOSSDA 대학생 데이터 시각화 공모전 2026
PPT 생성 스크립트 — 수상작 스타일 (빽빽한 데이터 중심 구성)
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Cm
from lxml import etree
import copy, os

BASE = '/Users/baiohelseu/Desktop/Project/kossda/'
OUT  = os.path.join(BASE, 'output/KOSSDA_2026_교권침해구조분석.pptx')

# ── 색상 팔레트 ──────────────────────────────────────────
NAVY       = RGBColor(0x0A, 0x24, 0x63)
NAVY_MID   = RGBColor(0x1B, 0x4F, 0x8A)
RED        = RGBColor(0xC0, 0x39, 0x2B)
TEAL       = RGBColor(0x00, 0x79, 0x7B)
ORANGE     = RGBColor(0xD3, 0x5A, 0x00)
LIGHT_BLUE = RGBColor(0xE3, 0xF0, 0xFB)
LIGHT_RED  = RGBColor(0xFD, 0xED, 0xEB)
LIGHT_TEAL = RGBColor(0xE0, 0xF4, 0xF4)
GRAY_BG    = RGBColor(0xF4, 0xF6, 0xF8)
GRAY_ROW   = RGBColor(0xEA, 0xEC, 0xEF)
TEXT_DARK  = RGBColor(0x12, 0x14, 0x20)
TEXT_MED   = RGBColor(0x4A, 0x4D, 0x5E)
TEXT_LIGHT = RGBColor(0x8A, 0x8D, 0x9E)
WHITE_WARM = RGBColor(0xFA, 0xFA, 0xFB)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)

KR = "Apple SD Gothic Neo"   # macOS 한글
EN = "Calibri"

# ── 슬라이드 사이즈: 16:9 ───────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # 완전 빈 레이아웃


# ═══════════════════════════════════════════════════════════
# 헬퍼 함수
# ═══════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=Pt(0)):
    """사각형 도형 추가"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             font_size=11, bold=False, italic=False,
             color=TEXT_DARK, align=PP_ALIGN.LEFT,
             font_name=KR, wrap=True, v_anchor=None):
    """텍스트 박스 추가"""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    if v_anchor:
        tf.vertical_anchor = v_anchor
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_text_box(slide, lines, x, y, w, h,
                 font_size=10, bold_first=False,
                 fill=None, border=None, color=TEXT_DARK,
                 align=PP_ALIGN.LEFT, line_spacing=None):
    """
    lines: list of (text, size, bold, color) or str
    """
    if fill:
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid(); shape.fill.fore_color.rgb = fill
        if border:
            shape.line.color.rgb = border; shape.line.width = Pt(1)
        else:
            shape.line.fill.background()
        tf = shape.text_frame
    else:
        txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, sz, bd, col = line, font_size, (bold_first and i==0), color
        else:
            txt, sz, bd, col = line
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = align
        if line_spacing:
            p.line_spacing = Pt(line_spacing)
        run = p.add_run()
        run.text = txt
        run.font.name = KR
        run.font.size = Pt(sz)
        run.font.bold = bd
        run.font.color.rgb = col


def add_image(slide, img_path, x, y, w, h=None):
    """이미지 삽입"""
    if h:
        slide.shapes.add_picture(img_path, Inches(x), Inches(y), Inches(w), Inches(h))
    else:
        slide.shapes.add_picture(img_path, Inches(x), Inches(y), width=Inches(w))


def add_table(slide, data, col_widths, x, y, total_w, row_h=Inches(0.28),
              header_fill=NAVY, header_text=WHITE,
              alt_fill=GRAY_ROW, font_size=8.5):
    """테이블 추가 (data: list of list of str)"""
    rows, cols = len(data), len(data[0])
    tbl = slide.shapes.add_table(rows, cols, Inches(x), Inches(y),
                                  Inches(total_w), row_h * rows).table
    tbl.first_row = True

    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = Inches(cw)

    for ri, row in enumerate(data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.text = cell_text
            # 폰트
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run() if not p.runs else p.runs[0]
            if not p.runs:
                run = p.add_run()
                run.text = cell_text
                p.runs[0].text = ''
            else:
                p.runs[0].font.name = KR
                p.runs[0].font.size = Pt(font_size)
                p.runs[0].font.bold = (ri == 0)
                if ri == 0:
                    p.runs[0].font.color.rgb = header_text
            # 헤더 배경
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            solidFill = etree.SubElement(tcPr, qn('a:solidFill'))
            srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'))
            if ri == 0:
                srgbClr.set('val', str(header_fill))
            elif ri % 2 == 0:
                srgbClr.set('val', str(alt_fill))
            else:
                srgbClr.set('val', 'FFFFFF')

    # 폰트 재설정 (테이블 셀의 기본 폰트 덮어쓰기)
    for ri, row in enumerate(data):
        for ci, cell_text in enumerate(row):
            cell = tbl.cell(ri, ci)
            p = cell.text_frame.paragraphs[0]
            p.clear()
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell_text
            run.font.name = KR
            run.font.size = Pt(font_size)
            run.font.bold = (ri == 0)
            if ri == 0:
                run.font.color.rgb = header_text
            else:
                run.font.color.rgb = TEXT_DARK
    return tbl


def slide_header(slide, section_num, section_title, accent=NAVY):
    """슬라이드 상단 헤더 바 추가"""
    # 배경 바
    add_rect(slide, 0, 0, 13.33, 0.62, fill=accent)
    # 섹션 번호 배지
    add_rect(slide, 0.2, 0.08, 0.55, 0.46, fill=WHITE)
    add_text(slide, section_num, 0.2, 0.10, 0.55, 0.42,
             font_size=16, bold=True, color=accent, align=PP_ALIGN.CENTER)
    # 제목
    add_text(slide, section_title, 0.85, 0.08, 12.0, 0.46,
             font_size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT)


def stat_box(slide, value, label, x, y, w=1.8, h=1.0,
             fill=LIGHT_BLUE, val_color=NAVY, lbl_color=TEXT_MED,
             val_size=28, lbl_size=9):
    """큰 숫자 통계 박스"""
    add_rect(slide, x, y, w, h, fill=fill)
    add_text(slide, value, x+0.05, y+0.05, w-0.1, h*0.55,
             font_size=val_size, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x+0.05, y+h*0.55, w-0.1, h*0.42,
             font_size=lbl_size, bold=False, color=lbl_color, align=PP_ALIGN.CENTER)


def callout_box(slide, title, body, x, y, w, h,
                fill=LIGHT_BLUE, border=NAVY_MID,
                title_color=NAVY, body_color=TEXT_DARK,
                title_size=10, body_size=9):
    """핵심 메시지 콜아웃 박스"""
    add_rect(slide, x, y, w, h, fill=fill, line_color=border, line_w=Pt(1.5))
    if title:
        add_text(slide, title, x+0.12, y+0.07, w-0.2, 0.28,
                 font_size=title_size, bold=True, color=title_color)
    add_text(slide, body, x+0.12, y+(0.32 if title else 0.1), w-0.2,
             h-(0.45 if title else 0.18),
             font_size=body_size, color=body_color, wrap=True)


def section_divider(slide, y=0.70, color=GRAY_BG):
    """가로 구분선"""
    add_rect(slide, 0.25, y, 12.83, 0.03, fill=color)


# ═══════════════════════════════════════════════════════════
# SLIDE 00 — 표지
# ═══════════════════════════════════════════════════════════
def make_cover(prs):
    sld = prs.slides.add_slide(BLANK)

    # 전체 네이비 배경
    add_rect(sld, 0, 0, 13.33, 7.5, fill=NAVY)

    # 상단 포인트 바
    add_rect(sld, 0, 0, 13.33, 0.18, fill=RED)

    # 하단 포인트 바
    add_rect(sld, 0, 7.2, 13.33, 0.30, fill=RGBColor(0x07, 0x15, 0x38))

    # KOSSDA 배지
    add_rect(sld, 0.5, 0.4, 2.8, 0.42, fill=NAVY_MID)
    add_text(sld, 'KOSSDA 대학생 데이터 시각화 공모전 2026',
             0.52, 0.42, 2.76, 0.38, font_size=8.5, bold=True,
             color=RGBColor(0xA0, 0xC4, 0xFF), align=PP_ALIGN.CENTER)

    # 메인 제목
    add_text(sld,
             '"인권교육 참여율이 20%p 증가했는데도\n교권침해는 왜 줄지 않는가?"',
             0.5, 1.05, 12.33, 2.4,
             font_size=38, bold=True, color=WHITE_WARM, align=PP_ALIGN.LEFT)

    # 부제목
    add_rect(sld, 0.5, 3.55, 8.0, 0.04, fill=RED)
    add_text(sld, '— 제도 공백의 구조, 취약 집단, 그리고 맞춤 처방',
             0.5, 3.70, 12.0, 0.55,
             font_size=19, bold=False, color=RGBColor(0xB0, 0xD0, 0xFF),
             align=PP_ALIGN.LEFT)

    # 구분선
    add_rect(sld, 0.5, 4.35, 12.33, 0.025, fill=RGBColor(0x30, 0x50, 0x90))

    # 핵심 수치 박스 (4개)
    stat_items = [
        ('85.6%', '교권침해 경험률', LIGHT_BLUE, NAVY),
        ('16.9%', '자살사고율', LIGHT_RED, RED),
        ('OR=2.57', '침해→자살사고', LIGHT_TEAL, TEAL),
        ('47.3%', '교사 제도적 긴장\n인식 비율', RGBColor(0xFF,0xF3,0xE0), ORANGE),
    ]
    for i, (val, lbl, fc, vc) in enumerate(stat_items):
        sx = 0.5 + i * 3.21
        add_rect(sld, sx, 4.5, 3.0, 1.0, fill=fc)
        add_text(sld, val, sx+0.1, 4.55, 2.8, 0.5,
                 font_size=26, bold=True, color=vc, align=PP_ALIGN.CENTER)
        add_text(sld, lbl, sx+0.1, 5.08, 2.8, 0.38,
                 font_size=9, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # 데이터 출처
    add_text(sld, '주요 데이터: 교원 인권상황 실태조사 2024 (n=10,888) · 초중등 교원 인권교육 실태조사 2021 (n=9,553)',
             0.5, 5.7, 12.33, 0.35, font_size=8.5,
             color=RGBColor(0x70, 0x90, 0xC0), align=PP_ALIGN.LEFT)

    # 분석 방법
    add_text(sld, '분석 기법: K-means · SEM · SHAP(RandomForest) · Bootstrap 매개효과 · PAF 시뮬레이션 · Two-Track 로지스틱 회귀',
             0.5, 6.05, 12.33, 0.35, font_size=8.5,
             color=RGBColor(0x70, 0x90, 0xC0), align=PP_ALIGN.LEFT)

    # 하단 선
    add_text(sld, '※ 본 분석은 KOSSDA 공개 데이터를 활용한 학술 연구 목적 분석임',
             0.5, 7.22, 12.33, 0.25, font_size=7.5,
             color=RGBColor(0x60, 0x80, 0xB0), align=PP_ALIGN.LEFT)


# ═══════════════════════════════════════════════════════════
# SLIDE 01 — 문제 제기
# ═══════════════════════════════════════════════════════════
def make_slide01(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '01', '문제 제기: 교사가 병들고 있다')

    # 핵심 메시지
    callout_box(sld, '핵심 메시지',
                '"교사 10명 중 8.5명이 침해를 경험하며, 이 중 16.9%는 자살을 생각한다"\n'
                '→ 교권침해는 개인 갈등이 아닌 교사 집단의 구조적 위기다.',
                0.25, 0.72, 12.83, 0.80,
                fill=LIGHT_BLUE, border=NAVY_MID,
                title_size=10, body_size=10)

    # 5개 대형 통계 박스
    stats = [
        ('85.6%', '교권침해 경험률', LIGHT_BLUE, NAVY),
        ('59.6%', '보호자 민원 경험', RGBColor(0xEA,0xF5,0xEA), TEAL),
        ('46.3%', '이직 고려율', RGBColor(0xFF,0xF3,0xE0), ORANGE),
        ('16.9%', '자살사고율', LIGHT_RED, RED),
        ('42.8%', '정신건강 나쁨', GRAY_BG, TEXT_MED),
    ]
    for i, (val, lbl, fc, vc) in enumerate(stats):
        sx = 0.25 + i * 2.57
        stat_box(sld, val, lbl, sx, 1.65, w=2.45, h=1.05,
                 fill=fc, val_color=vc, val_size=30, lbl_size=9.5)

    # 데이터 테이블
    table_data = [
        ['수치', '값', '출처 변수', '의미'],
        ['교권침해 경험률', '85.6%', 'B2_1+B3_1+B5_1', '학생·보호자·관리자 침해 합산'],
        ['보호자 민원 경험', '59.6%', 'B3_1', '보호자로부터 침해 경험'],
        ['이직 고려율', '46.3%', 'A4', '현 직장 이직을 고려함'],
        ['자살사고율', '16.9%', 'C4', '최근 자살을 생각해본 경험'],
        ['정신건강 나쁨', '42.8%', 'C1 역코딩', '정신건강 자가 평가 하위'],
    ]
    add_table(sld, table_data, [3.2, 1.2, 2.0, 5.0],
              0.25, 2.88, 12.83, row_h=Inches(0.29), font_size=9)

    # 스토리 본문
    add_text(sld,
             '교권침해는 더 이상 개인 갈등이 아니다. 2024년 전국 교원 조사(n=10,888) 결과 교권침해 경험률 85.6%,'
             ' 이직 고려 46.3%, 자살사고 16.9%로 나타났다. 이는 교사 집단 전체가 구조적 위기에 놓여 있음을 보여준다.',
             0.25, 4.83, 8.0, 0.70,
             font_size=9.5, color=TEXT_DARK, wrap=True)

    # 연구 질문 박스
    add_rect(sld, 8.45, 4.72, 4.63, 1.48, fill=NAVY)
    add_text(sld, '핵심 연구 질문',
             8.57, 4.80, 4.39, 0.32,
             font_size=9, bold=True, color=RGBColor(0xA0,0xC4,0xFF))
    add_text(sld,
             '"인권교육 참여율이 20%p 증가했는데도\n교권침해는 왜 줄지 않는가?\n— 제도 공백의 구조를 규명한다"',
             8.57, 5.14, 4.39, 1.0,
             font_size=11, bold=True, color=WHITE_WARM, wrap=True)

    # 출처 및 방법
    add_text(sld, '출처: D1 교원 인권상황 실태조사(2024), n=10,888',
             0.25, 5.62, 8.0, 0.28,
             font_size=8, color=TEXT_LIGHT, italic=True)

    # 하단 이미지 (final_story 일부)
    img_path = BASE + 'output/partA_teacher_rights_survey.png'
    if os.path.exists(img_path):
        add_image(sld, img_path, 0.25, 5.98, 12.83, 1.30)


# ═══════════════════════════════════════════════════════════
# SLIDE 02 — 데이터와 분석 방법
# ═══════════════════════════════════════════════════════════
def make_slide02(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '02', '데이터와 분석 방법: Two-Track 전략으로 역인과를 해소한다', accent=NAVY_MID)

    # 핵심 메시지 배너
    callout_box(sld, None,
                '5개 데이터 · 9가지 통계 기법 · Two-Track 설계 — D2(제도 필요도)의 역인과 문제를 구조적으로 해소',
                0.25, 0.72, 12.83, 0.48, fill=GRAY_BG, border=NAVY_MID,
                body_size=10.5, body_color=NAVY)

    # 왼쪽: 데이터 테이블
    add_text(sld, '분석 데이터', 0.25, 1.30, 6.2, 0.30, font_size=11, bold=True, color=NAVY)
    data_table = [
        ['데이터', 'n', '역할'],
        ['교원 인권상황 실태조사 (2024)', '10,888', '주 분석 데이터'],
        ['초·중등 교원 인권교육 실태조사 (2021)', '9,553', '교육 효과성 + 시계열'],
        ['한국교총 교권상담 실적 (2018)', '501건', '신고 추이 외부 보조'],
        ['신도시 현황 + 순이동인구', '32개 단지', '지역 구조 변수'],
        ['학원·교습소 현황 (시군구별)', '시군구', '교육열 프록시'],
    ]
    add_table(sld, data_table, [4.8, 1.3, 2.7],
              0.25, 1.62, 6.8, row_h=Inches(0.295), font_size=9)

    # 오른쪽: 분석 기법
    add_text(sld, '적용 기법 (9종)', 7.2, 1.30, 5.88, 0.30, font_size=11, bold=True, color=NAVY)
    methods = [
        ['기법', '적용 목적'],
        ['K-means (k=3)', '취약 집단 분류 (Silhouette=0.404)'],
        ['SEM (semopy)', '구조적 인과 경로 분석'],
        ['SHAP (RandomForest)', '제도 변수 중요도 해석'],
        ['Brunner-Munzel', '비모수 집단 비교 (η²)'],
        ['Bootstrap 매개효과', '간접효과 95%CI (n=2,000)'],
        ['5-fold CV', '모델 예측력 교차 검증'],
        ['VIF', '다중공선성 진단'],
        ['PAF 시뮬레이션', '정책 효과 추정 (역학 표준)'],
        ['Two-Track 로지스틱', '역인과 해소 핵심 전략'],
    ]
    add_table(sld, methods, [2.8, 4.3],
              7.2, 1.62, 5.88, row_h=Inches(0.265), font_size=9)

    # Two-Track 설명 박스
    add_text(sld, 'Two-Track 설계 (역인과 해소 핵심)', 0.25, 3.90, 12.83, 0.30,
             font_size=11, bold=True, color=RED)
    track_table = [
        ['Track', '방법', '역인과 여부', '역할'],
        ['A', '구조변수(학교급·규모·고용형태) → 침해 로지스틱 OR', '없음 ✓', '인과 추론 가능 취약 집단 식별'],
        ['B', '침해집단 vs 비침해집단 D2 수요 격차 분석', '강점으로 전환 ✓', '당사자 기반 정책 우선순위 도출'],
    ]
    add_table(sld, track_table, [0.9, 5.5, 2.0, 4.2],
              0.25, 4.22, 12.83, row_h=Inches(0.35),
              header_fill=RED, font_size=9.5)

    # D2 내생성 설명
    add_rect(sld, 0.25, 5.10, 12.83, 1.20, fill=LIGHT_RED)
    add_text(sld, 'D2(제도 필요도) 내생성 처리 원칙',
             0.37, 5.15, 8.0, 0.30, font_size=10, bold=True, color=RED)
    add_text(sld,
             'D2는 침해 경험 이후 묻는 인식 변수이므로 "제도 없어서 침해 발생"의 역인과 가능성이 있음.\n'
             '→ Track A에서 D2 완전 제외 (구조 변수만 사용) · Track B에서 "당사자 수요 측정 도구"로만 사용 (인과 주장 없음)\n'
             '→ PAF 시뮬레이션의 OR은 수요 기반 상한 추정치로 해석',
             0.37, 5.46, 12.5, 0.80, font_size=9.5, color=TEXT_DARK)

    # 하단 적합도
    add_text(sld,
             'SEM 모델 적합도: CFI=0.885 (경계), RMSEA=0.060 (수용) '
             '| RandomForest 5-fold CV AUC: 이직고려 0.71 / 자살사고 0.73',
             0.25, 6.42, 12.83, 0.35,
             font_size=8.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 03 — 발견 1: 인권교육 역설
# ═══════════════════════════════════════════════════════════
def make_slide03(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '03', '발견 1: 인권교육 역설 [HOOK] — 교육량이 아닌 제도의 문제다', accent=TEAL)

    callout_box(sld, '핵심 메시지',
                '교육 참여율 +20.8%p ↑, 그러나 효과 만족도 3.75/5 · 교사 47.3%가 "학생인권↑=교권↓" 인식\n'
                '→ 인식·태도 교육만으로는 구조적 침해를 막을 수 없음을 3가지 지표로 실증',
                0.25, 0.72, 12.83, 0.75,
                fill=LIGHT_TEAL, border=TEAL, title_size=10, body_size=10)

    # 시계열 비교 (3개 큰 숫자)
    paradox_items = [
        ('67.4%', '2021년 인권교육\n참여율 (D3)', LIGHT_BLUE, NAVY),
        ('88.2%', '2024년 인권교육\n참여율 (D1)', NAVY, WHITE),
        ('+20.8%p', '3년간 증가폭', LIGHT_TEAL, TEAL),
    ]
    for i, (val, lbl, fc, vc) in enumerate(paradox_items):
        sx = 0.25 + i * 2.95
        stat_box(sld, val, lbl, sx, 1.57, w=2.80, h=0.98,
                 fill=fc, val_color=vc, val_size=28, lbl_size=9)

    # 화살표 텍스트
    add_text(sld, '→', 3.15, 1.87, 0.5, 0.5, font_size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    add_text(sld, '→', 6.10, 1.87, 0.5, 0.5, font_size=28, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

    # BUT 박스
    add_rect(sld, 9.2, 1.57, 3.88, 0.98, fill=LIGHT_RED)
    add_text(sld, 'BUT', 9.3, 1.60, 3.68, 0.35, font_size=22, bold=True, color=RED, align=PP_ALIGN.CENTER)
    add_text(sld, '교권침해 85.6% 여전\n자살사고 16.9%', 9.3, 1.97, 3.68, 0.55,
             font_size=10, color=RED, align=PP_ALIGN.CENTER)

    # 왼쪽: 상세 비교 테이블
    add_text(sld, '시계열 비교 — D3(2021) × D1(2024)', 0.25, 2.68, 6.5, 0.28,
             font_size=10.5, bold=True, color=TEAL)
    compare_table = [
        ['지표', '2021 (D3)', '2024 (D1)', '변화 방향'],
        ['인권교육 참여율', '67.4%', '88.2%', '+20.8%p ↑'],
        ['교육 효과 만족도 (B2_6)', '3.75/5.0', '—', '보통 이하 39.1%'],
        ['인권의식 향상 응답 (B2_9≥4)', '62.3%', '—', '낮은 자기 효능감'],
        ['"학생인권↑→교권↓" 동의 (B5_7≥4)', '47.3%', '—', '제도적 긴장 지속'],
        ['교권침해 경험률', '측정 없음', '85.6%', '여전히 높음'],
        ['자살사고율', '측정 없음', '16.9%', '—'],
    ]
    add_table(sld, compare_table, [3.6, 1.5, 1.5, 2.2],
              0.25, 2.98, 8.85, row_h=Inches(0.285), font_size=9)

    # 오른쪽: D3 교육 효과 그림
    img_path = BASE + 'output/two_track_analysis.png'
    if os.path.exists(img_path):
        # two_track_analysis의 교육효과 패널 (상단 오른쪽)만 보여주기 위해 전체 삽입
        add_image(sld, img_path, 9.25, 2.68, 3.83, 2.85)

    # 3줄 논거 박스
    add_rect(sld, 0.25, 5.62, 12.83, 1.62, fill=LIGHT_BLUE)
    add_text(sld, '논거 요약 (3중 증거)',
             0.37, 5.68, 5.0, 0.30, font_size=10, bold=True, color=NAVY)
    evidence = [
        '① 참여율 88.2%에도 교육 만족도 3.75/5 → 교육 자체의 효과성이 낮음',
        '② 교사 47.3%: "학생인권 확대가 교권을 축소한다" → 구조적 위협 인식 지속',
        '③ 인권의식 향상 응답 62.3%에 불과 → 인식 변화가 곧 행동·구조 변화로 연결되지 않음',
        '→ 결론: 인식 개선 교육만으로는 구조적 침해를 막을 수 없다. 제도 설계의 문제다.',
    ]
    for i, ev in enumerate(evidence):
        col = RED if i == 3 else TEXT_DARK
        bd = (i == 3)
        add_text(sld, ev, 0.37, 5.98 + i*0.29, 12.5, 0.26,
                 font_size=9.5, color=col, bold=bd)

    add_text(sld, '출처: D3 초·중등 교원 인권교육 실태조사(2021) B1·B2_6·B2_9·B5_7 / D1 교원 인권상황 실태조사(2024)',
             0.25, 7.28, 12.83, 0.20, font_size=7.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 04 — 발견 2: Track A 구조→침해
# ═══════════════════════════════════════════════════════════
def make_slide04(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '04', '발견 2: 구조가 침해를 만든다 — Track A (역인과 없는 인과 추론)', accent=RGBColor(0x1B,0x5E,0x20))

    callout_box(sld, '핵심 메시지',
                '초등학교 OR=1.79***, 학교규모 1단계↑ OR=1.29*** — 구조 변수만으로 취약 집단 식별 가능\n'
                '침해는 이직(OR=2.17***)·자살사고(OR=2.57***)를 직접 유발한다 (구조 통제 후)',
                0.25, 0.72, 12.83, 0.75,
                fill=RGBColor(0xE8,0xF5,0xE9), border=RGBColor(0x1B,0x5E,0x20),
                title_size=10, body_size=10,
                title_color=RGBColor(0x1B,0x5E,0x20), body_color=TEXT_DARK)

    # 왼쪽: Track A 구조 OR 테이블
    add_text(sld, 'Track A: 구조 변수 → 침해 OR (n=10,888, 구조 변수 동시 투입)',
             0.25, 1.57, 7.9, 0.28, font_size=10.5, bold=True, color=RGBColor(0x1B,0x5E,0x20))
    or_table = [
        ['구조 변수', '보호자침해 OR', '학생침해 OR', '관리자침해 OR', '해석'],
        ['초등학교 (vs 기준)', '1.470***', '1.790***', '0.550***', '초등 보호자/학생침해 취약'],
        ['중학교 (vs 기준)', '0.885 n.s.', '1.176 n.s.', '0.629***', '중학교 상대적 보통'],
        ['고등학교 (vs 기준)', '0.619***', '0.660***', '0.683***', '고등 상대적 보호적'],
        ['학교규모 (1단계↑)', '1.080***', '1.289***', '0.939**', '학교 클수록 학생침해↑'],
        ['사립 (vs 공립)', '0.800 n.s.', '0.951 n.s.', '1.195 n.s.', '설립 형태 영향 미미'],
        ['기간제 교원', '0.689**', '0.668**', '0.782 n.s.', '※ 신고 억제 가능성 주의'],
        ['남성 (vs 여성)', '0.817***', '0.663***', '0.908 n.s.', '여성 교원 상대적 취약'],
    ]
    add_table(sld, or_table, [2.5, 1.5, 1.5, 1.5, 2.7],
              0.25, 1.88, 9.7, row_h=Inches(0.27), font_size=8.5)

    # 오른쪽: 침해→결과 OR
    add_text(sld, '침해 → 결과 직접효과 (구조 통제 후)',
             10.1, 1.57, 3.0, 0.28, font_size=10, bold=True, color=RED)
    result_table = [
        ['결과 변수', 'OR', '95%CI'],
        ['이직 고려', '2.174***', '[2.07, 2.28]'],
        ['자살사고', '2.569***', '[2.40, 2.75]'],
    ]
    add_table(sld, result_table, [1.6, 0.8, 1.35],
              10.1, 1.88, 3.0, row_h=Inches(0.32),
              header_fill=RED, font_size=9.5)

    # 번아웃 재검토 박스
    add_rect(sld, 0.25, 3.88, 12.83, 0.95, fill=LIGHT_RED)
    add_text(sld, '번아웃 매개 경로 재검토 (데이터 기반 수정)',
             0.37, 3.93, 8.0, 0.28, font_size=10, bold=True, color=RED)
    add_text(sld,
             '실측: 침해집단의 번아웃(C3) 점수가 비침해집단보다 낮음 (보호자침해: 침해 1.90 vs 비침해 2.37, Δ=-0.47***)\n'
             '→ 번아웃과 침해는 구조적 요인에 의해 각각 독립 예측되는 별개 경로 — 매개 주장 성립 불안정\n'
             '→ 인과 주장을 "침해→이직/자살 직접효과"로 집중 (Track A의 OR=2.17·2.57***로 충분히 지지)',
             0.37, 4.23, 12.5, 0.57, font_size=9.5, color=TEXT_DARK)

    # 기간제 주의 박스
    add_rect(sld, 0.25, 4.95, 6.2, 0.82, fill=RGBColor(0xFF,0xF3,0xE0))
    add_text(sld, '기간제 OR 해석 주의',
             0.37, 5.00, 5.0, 0.25, font_size=9.5, bold=True, color=ORANGE)
    add_text(sld,
             '기간제 OR이 낮은 것은 보호 효과가 아닌 "계약 유지 압박으로 인한 침해 신고 억제" 가능성.\n'
             '기간제의 실제 침해 규모가 과소 추정될 수 있는 측정 한계',
             0.37, 5.27, 5.8, 0.48, font_size=9, color=TEXT_DARK)

    # Two-Track 이미지 (하단)
    img_path = BASE + 'output/two_track_analysis.png'
    if os.path.exists(img_path):
        add_image(sld, img_path, 6.7, 4.90, 6.38, 2.38)

    add_text(sld,
             '출처: D1 교원 인권상황 실태조사(2024) | 분석: analysis_two_track.py | '
             '기준집단: 유치원+특수+기타(SQ2=1,6,7) | *** p<.001  ** p<.01  * p<.05',
             0.25, 7.28, 12.83, 0.20, font_size=7.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 05 — 발견 3: Two-Track 수렴
# ═══════════════════════════════════════════════════════════
def make_slide05(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '05', '발견 3: Track A × Track B 수렴 — 수렴 타당도로 확정하는 정책 우선순위')

    callout_box(sld, '핵심 메시지',
                '구조적 취약 집단(Track A)과 당사자 정책 수요(Track B)가 같은 제도를 지목 → 수렴 타당도 확보\n'
                '전 교사 95%+가 모든 정책 필요 인식 + 침해집단 더욱 절실히 요구 (수요격차 *** 유의)',
                0.25, 0.72, 12.83, 0.75, fill=LIGHT_BLUE, border=NAVY,
                title_size=10, body_size=10)

    # Track B 수요 격차 테이블
    add_text(sld, 'Track B: D2 정책 수요 격차 (역인과를 "당사자 수요 근거"로 전환)',
             0.25, 1.57, 8.0, 0.28, font_size=10.5, bold=True, color=NAVY)
    demand_table = [
        ['정책 제도', '전체 수요(≥4 비율)', '보호자침해집단 Δ', '학생침해집단 Δ', '관리자침해집단 Δ'],
        ['악성민원 법적패널티 (D2_10)', '99.0%', '+0.088***', '+0.109***', '+0.025***'],
        ['수업방해 분리시스템 (D2_5)', '97.9%', '+0.075***', '+0.128***', '+0.015***'],
        ['아동학대 신고보호 (D2_2)', '99.0%', '+0.072***', '+0.082***', '+0.018**'],
        ['과밀학급 해소 (D2_16)', '96.8%', '+0.043***', '+0.078***', '+0.045***'],
    ]
    add_table(sld, demand_table, [3.5, 1.8, 1.9, 1.9, 2.0],
              0.25, 1.88, 11.1, row_h=Inches(0.28), font_size=9)

    # SHAP 테이블
    add_text(sld, 'SHAP 중요도 (RandomForest, n=10,888, n_estimators=300)',
             0.25, 3.30, 5.5, 0.28, font_size=10, bold=True, color=NAVY_MID)
    shap_table = [
        ['순위', '변수', 'mean|SHAP|', 'Track 일치'],
        ['1위', '학교 규모', '0.035', 'Track A ✓'],
        ['2위', '초등학교 여부', '0.033', 'Track A ✓'],
        ['3위', '수업방해 분리시스템', '0.023', 'Track B ✓'],
        ['4위', '악성민원 법적패널티', '0.018', 'Track B ✓'],
    ]
    add_table(sld, shap_table, [0.8, 2.8, 1.3, 1.3],
              0.25, 3.60, 6.2, row_h=Inches(0.28), font_size=9)

    # 수렴 결론 테이블
    add_text(sld, 'Track A × Track B 수렴 — 정책 우선순위 확정',
             0.25, 5.02, 8.0, 0.28, font_size=10.5, bold=True, color=RED)
    convergence_table = [
        ['침해 유형', 'Track A (구조 취약성)', 'Track B (수요 격차)', '수렴 결론'],
        ['보호자침해 (59.6%)', '초등 OR=1.47 / 대형 OR=1.08', 'Δ악성민원=+0.088 최대', '악성민원패널티 1순위'],
        ['학생침해 (78.5%)', '초등 OR=1.79 / 대형 OR=1.29', 'Δ수업방해=+0.128 최대', '수업방해분리 1순위'],
        ['관리자침해 (23.6%)', 'D2 모두 비유의 (n=847)', '격차 미약 (측정 공백)', '신규 제도 필요 지적'],
    ]
    add_table(sld, convergence_table, [2.0, 3.2, 2.5, 2.8],
              0.25, 5.32, 10.5, row_h=Inches(0.30),
              header_fill=RED, font_size=9)

    # 오른쪽 이미지
    img_path = BASE + 'output/two_track_analysis.png'
    if os.path.exists(img_path):
        add_image(sld, img_path, 11.38, 3.30, 1.70, 3.92)

    # D2 해석 원칙
    add_rect(sld, 0.25, 6.65, 10.9, 0.60, fill=LIGHT_BLUE)
    add_text(sld,
             'D2 해석 원칙: D2는 침해 후 측정된 인식 변수 → 인과 추론(X) · 정책 수요 진단(O)\n'
             '"침해 경험자가 절실히 요구하는 정책 = 현장 기반 최우선 입법 타깃" — 역인과를 강점으로 전환',
             0.37, 6.70, 10.7, 0.52, font_size=9.5, color=NAVY)

    add_text(sld, '출처: D1(2024) | 분석: analysis_two_track.py, analysis_factor_impact.py',
             0.25, 7.28, 12.83, 0.20, font_size=7.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 06 — 발견 4: K-means 클러스터링
# ═══════════════════════════════════════════════════════════
def make_slide06(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '06', '발견 4: 취약 집단 분류 — 학교마다 다른 침해 구조 (K-means, k=3)')

    callout_box(sld, '핵심 메시지',
                'K-means 3집단: 각 집단이 완전히 다른 침해 구조를 가짐 — Track A 구조 OR 패턴과 일치 (수렴 검증)\n'
                'Silhouette Score=0.404 | D2(제도 필요도) 완전 제외 후 순수 구조·침해 프로파일로 분류',
                0.25, 0.72, 12.83, 0.75, fill=LIGHT_BLUE, border=NAVY,
                title_size=10, body_size=10)

    # 클러스터 3개 박스
    clusters = [
        ('C0: 고등·학생중심형', '1,286명 (11.8%)', '고등학교 86% | 공립 80%\n학생침해 72%', NAVY, WHITE),
        ('C1: 중학교·보호자민원형', '8,755명 (80.4%)', '중학교 86% | 공립 100% | 대형\n보호자침해 62%', TEAL, WHITE),
        ('C2: 초등·관리자압박형', '847명 (7.8%)', '초등 100% | 소형(SQ5=1.6)\n관리자침해 43%', RED, WHITE),
    ]
    for i, (title, n, desc, fc, tc) in enumerate(clusters):
        sx = 0.25 + i * 4.36
        add_rect(sld, sx, 1.57, 4.15, 0.35, fill=fc)
        add_text(sld, title, sx+0.08, 1.60, 3.99, 0.30,
                 font_size=11, bold=True, color=tc)
        add_rect(sld, sx, 1.92, 4.15, 0.28, fill=RGBColor(0xF0,0xF0,0xF5))
        add_text(sld, n, sx+0.08, 1.93, 3.99, 0.25, font_size=9, color=TEXT_MED)
        add_rect(sld, sx, 2.20, 4.15, 0.52, fill=WHITE)
        add_text(sld, desc, sx+0.08, 2.22, 3.99, 0.48, font_size=9, color=TEXT_DARK)

    # 클러스터 내 OR 테이블
    add_text(sld, '클러스터 내 로지스틱 OR (구조 변수 통제 후)',
             0.25, 2.85, 8.0, 0.28, font_size=10.5, bold=True, color=NAVY)
    cluster_or = [
        ['제도', 'C0 고등(n=1,286)', 'C1 중학교(n=8,755)', 'C2 초등(n=847)'],
        ['악성민원 법적패널티', 'OR=2.018***', 'OR=2.380***', 'OR=1.757**'],
        ['수업방해 분리시스템', 'OR=1.505***', 'OR=1.624***', 'OR=1.224 n.s.'],
        ['아동학대 신고보호', 'OR=1.610***', 'OR=1.883***', 'OR=1.398 n.s.'],
        ['과밀학급 해소', 'OR=1.298***', 'OR=1.335***', 'OR=1.087 n.s.'],
    ]
    add_table(sld, cluster_or, [3.5, 2.5, 2.5, 2.2],
              0.25, 3.15, 10.7, row_h=Inches(0.285), font_size=9)

    # Track A 연계
    add_rect(sld, 0.25, 4.80, 12.83, 0.60, fill=RGBColor(0xE8,0xF5,0xE9))
    add_text(sld, 'Track A 연계 검증:',
             0.37, 4.84, 2.5, 0.25, font_size=9.5, bold=True, color=RGBColor(0x1B,0x5E,0x20))
    add_text(sld,
             'C1(중학교·대형) ↔ Track A "학교규모 OR=1.29***" 일치 | C0(고등·학생침해) ↔ Track A "고등 OR=0.66 보호방향" 일치\n'
             'C2(초등·관리자) ↔ Track A "초등 OR=1.47-1.79***" 일치 | K-means가 Track A 구조 패턴을 독립 재현 → 수렴 타당도',
             2.9, 4.85, 10.0, 0.52, font_size=9.5, color=TEXT_DARK)

    # 이미지
    img_path = BASE + 'output/clustering_analysis.png'
    if os.path.exists(img_path):
        add_image(sld, img_path, 0.25, 5.50, 8.5, 1.75)

    img_path2 = BASE + 'output/policy_cards.png'
    if os.path.exists(img_path2):
        add_image(sld, img_path2, 8.9, 5.50, 4.18, 1.75)

    add_text(sld, '출처: D1(2024) | 분석: analysis_clustering.py | D2 제외, 순수 구조·침해 프로파일 클러스터링',
             0.25, 7.28, 12.83, 0.20, font_size=7.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 07 — 발견 5: 사각지대
# ═══════════════════════════════════════════════════════════
def make_slide07(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '07', '발견 5: 사각지대 — 초등 관리자압박형은 현행 입법 논의에서 빠져있다', accent=RED)

    callout_box(sld, '핵심 메시지',
                'C2 초등·관리자압박형(n=847): 관리자침해 43.2% 최고 — 그러나 현행 D2 4개 제도 모두 비유의\n'
                '→ "효과 없음"이 아니라 "측정 도구 부재" + "입법 논의의 사각지대" 2중 의미',
                0.25, 0.72, 12.83, 0.75,
                fill=LIGHT_RED, border=RED, title_size=10, body_size=10, title_color=RED)

    # 4개 지표 박스
    blind_stats = [
        ('43.2%', '관리자침해 비율\n(전체 평균 24.6%)', LIGHT_RED, RED),
        ('7.8%', 'C2 집단 비율\n(n=847 소표본)', GRAY_BG, TEXT_MED),
        ('4개 모두\nn.s.', 'D2 제도 효과\n(검정력 제한)', RGBColor(0xFF,0xFB,0xE6), ORANGE),
        ('1.6/5', '평균 학교 규모\n(소형 학교)', LIGHT_BLUE, NAVY_MID),
    ]
    for i, (val, lbl, fc, vc) in enumerate(blind_stats):
        sx = 0.25 + i * 3.22
        stat_box(sld, val, lbl, sx, 1.57, w=3.08, h=0.95,
                 fill=fc, val_color=vc, val_size=24, lbl_size=9)

    # 상세 비교 테이블
    add_text(sld, 'C2 초등·관리자압박형 상세 프로파일',
             0.25, 2.65, 7.0, 0.28, font_size=10.5, bold=True, color=RED)
    profile_table = [
        ['지표', 'C2 초등·관리자압박형', '전체 평균', '의미'],
        ['관리자침해 비율', '43.2%', '24.6%', 'C2가 1.76배 높음'],
        ['학교 규모 (SQ5)', '1.6/5 (소형)', '3.6/5', '소형 학교 집중'],
        ['공립 비율', '98.3%', '96.3%', '거의 전부 공립'],
        ['D2 악성민원패널티 OR', '1.757** (유의)', '(전체 OR=~2.38)', '상대적으로 낮음'],
        ['D2 수업방해분리 OR', '1.224 n.s.', '(전체 OR=~1.62)', '비유의 — 해당 없음'],
        ['D2 아동학대신고 OR', '1.398 n.s.', '(전체 OR=~1.88)', '비유의 — 해당 없음'],
    ]
    add_table(sld, profile_table, [2.8, 2.5, 1.8, 4.3],
              0.25, 2.95, 11.4, row_h=Inches(0.270), font_size=9)

    # 2중 의미 박스
    add_rect(sld, 0.25, 4.90, 8.0, 1.90, fill=LIGHT_RED)
    add_text(sld, '발견의 2중 의미',
             0.37, 4.96, 5.0, 0.28, font_size=10, bold=True, color=RED)
    meanings = [
        '① [입법 사각지대] 학부모·학생 중심 교권보호법 논의는 초등 소규모 학교 교사를 보호하지 못함',
        '② [측정 공백] n=847 소집단에서 검정력 부족 → 효과 부재가 아닌 검출 불가',
        '③ [전용 제도 부재] 관리자 갑질·부당 지시 대응 내부신고·보호 시스템이 D2에 없음',
        '→ 도구 개발 + 전용 입법이 선행 과제: 관리자 압박 내부신고 보호 시스템 신설',
    ]
    for i, m in enumerate(meanings):
        col = RED if i == 3 else TEXT_DARK
        add_text(sld, m, 0.37, 5.25 + i*0.36, 7.7, 0.33,
                 font_size=9.5, color=col, bold=(i == 3))

    # 오른쪽: 이미지
    img_path = BASE + 'output/clustering_analysis.png'
    if os.path.exists(img_path):
        add_image(sld, img_path, 8.5, 2.65, 4.58, 4.15)

    add_text(sld, '출처: D1(2024) | 분석: analysis_clustering.py | C2 n=847 (검정력 제한 명시)',
             0.25, 7.28, 12.83, 0.20, font_size=7.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 08 — 발견 6: 집단별 맞춤 처방
# ═══════════════════════════════════════════════════════════
def make_slide08(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '08', '발견 6: 집단별 맞춤 제도 처방 — 일률적 교권보호법의 한계')

    callout_box(sld, '핵심 메시지',
                '침해율 × Track B 수요격차 = 정책 우선순위 점수 — 같은 제도도 집단마다 우선순위 달라\n'
                'C2 초등은 현행 4개 제도 모두 해당 없음 → 신규 관리자압박 제도 별도 신설 필요',
                0.25, 0.72, 12.83, 0.75, fill=LIGHT_BLUE, border=NAVY,
                title_size=10, body_size=10)

    # 집단별 처방 카드 (3개)
    prescriptions = [
        ('C0', '고등·학생중심형', '수업방해\n분리시스템', '학생침해 72%\nTrack B Δ=+0.128***\nSHAP 3위 (0.023)', NAVY),
        ('C1', '중학교·보호자민원형', '악성민원\n법적패널티', '보호자침해 62%\nTrack B Δ=+0.088***\nOR=2.380 최고', TEAL),
        ('C2', '초등·관리자압박형', '관리자 압박\n내부신고 시스템', '관리자침해 43%\nD2 4개 모두 n.s.\n→ 신규 입법 필요', RED),
    ]
    for i, (code, name, policy, rationale, color) in enumerate(prescriptions):
        sx = 0.25 + i * 4.36
        add_rect(sld, sx, 1.57, 4.15, 0.38, fill=color)
        add_text(sld, f'{code}: {name}', sx+0.1, 1.60, 3.95, 0.33,
                 font_size=11, bold=True, color=WHITE)
        add_rect(sld, sx, 1.95, 4.15, 0.05, fill=RGBColor(0xE0,0xE0,0xE5))
        add_rect(sld, sx, 2.0, 4.15, 0.78, fill=LIGHT_BLUE if i < 2 else LIGHT_RED)
        add_text(sld, f'1순위 제도:\n{policy}', sx+0.1, 2.02, 3.95, 0.74,
                 font_size=12, bold=True,
                 color=color if i < 2 else RED, align=PP_ALIGN.CENTER)
        add_rect(sld, sx, 2.78, 4.15, 0.62, fill=GRAY_BG)
        add_text(sld, rationale, sx+0.1, 2.80, 3.95, 0.58,
                 font_size=9, color=TEXT_DARK)

    # 정책 로드맵 테이블
    add_text(sld, '정책 로드맵 (시급도 × 근거 강도)',
             0.25, 3.55, 10.0, 0.28, font_size=10.5, bold=True, color=NAVY)
    roadmap = [
        ['단계', '대상 집단', '정책 내용', '핵심 근거', '시급도'],
        ['즉시', 'C1 중학교 (80.4%)', '악성민원 법적패널티 입법', 'Track B Δ=+0.088*** / OR=2.380***', '★★★★★'],
        ['단기', 'C0 고등 + C1 중학교', '수업방해 학생분리 조항 신설', 'Track B Δ=+0.128*** / SHAP 3위', '★★★★'],
        ['중기', 'C1 중학교', '아동학대 신고 교원 보호 강화', 'Track B Δ=+0.072*** / OR=1.883***', '★★★'],
        ['신규', 'C2 초등 (사각지대)', '관리자압박 내부신고·보호 시스템', 'D2 4개 모두 n.s. → 도구 부재 확인', '★★★ (신규)'],
    ]
    add_table(sld, roadmap, [0.9, 2.5, 3.0, 4.2, 1.5],
              0.25, 3.85, 12.35, row_h=Inches(0.285),
              header_fill=NAVY, font_size=9)

    # 이미지
    img_path = BASE + 'output/policy_roadmap.png'
    if os.path.exists(img_path):
        add_image(sld, img_path, 0.25, 5.45, 8.5, 1.80)

    img_path2 = BASE + 'output/policy_cards.png'
    if os.path.exists(img_path2):
        add_image(sld, img_path2, 8.9, 5.45, 4.18, 1.80)

    add_text(sld, '출처: D1(2024) | 분석: analysis_clustering.py, analysis_two_track.py',
             0.25, 7.28, 12.83, 0.20, font_size=7.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 09 — PAF 시뮬레이션
# ═══════════════════════════════════════════════════════════
def make_slide09(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '09', '미래 전망: 제도 도입 시 침해율 감소 시뮬레이션 (PAF)', accent=ORANGE)

    callout_box(sld, '핵심 메시지',
                '악성민원 법적패널티 60% 효과 시: 보호자침해 59.6% → 38.3% (−21.3%p)\n'
                '※ PAF의 OR은 D2 수요 기반 — 보수적 상한 추정치로 사용 (인과 추론 한계 명시)',
                0.25, 0.72, 12.83, 0.75,
                fill=RGBColor(0xFF,0xF8,0xE1), border=ORANGE,
                title_size=10, body_size=10, title_color=ORANGE)

    # PAF 공식
    add_text(sld, 'PAF 산출 공식 (역학 표준 Population Attributable Fraction)',
             0.25, 1.57, 10.0, 0.28, font_size=10, bold=True, color=ORANGE)
    add_rect(sld, 0.25, 1.87, 12.83, 0.42, fill=RGBColor(0xFF,0xF8,0xE1))
    add_text(sld,
             'PAF = Pe × (OR − 1) / [1 + Pe × (OR − 1)]    '
             '| Pe: 제도 부재 노출 비율 (D2≥4 응답률) | OR: D2 기반 로지스틱 결과',
             0.37, 1.90, 12.5, 0.36, font_size=10, color=TEXT_DARK, align=PP_ALIGN.CENTER)

    # PAF 메인 테이블
    add_text(sld, 'PAF 시뮬레이션 결과 (3단계 효과 시나리오)',
             0.25, 2.42, 10.0, 0.28, font_size=10.5, bold=True, color=NAVY)
    paf_table = [
        ['제도', 'Pe', 'OR', 'PAF', '현재 침해율', '30% 효과', '60% 효과', '90% 효과'],
        ['악성민원 법적패널티', '99.0%', '2.487', '0.595', '59.6%', '49.0%', '38.3%', '27.7%'],
        ['수업방해 분리시스템', '97.9%', '1.710', '0.410', '78.5%', '68.8%', '59.2%', '49.5%'],
        ['아동학대 신고보호', '99.0%', '1.978', '0.492', '59.6%', '50.8%', '42.0%', '33.2%'],
    ]
    add_table(sld, paf_table, [2.8, 0.75, 0.75, 0.75, 1.25, 1.15, 1.15, 1.15],
              0.25, 2.73, 9.75, row_h=Inches(0.305),
              header_fill=ORANGE, font_size=9.5)

    # 60% 효과 강조 박스 (오른쪽)
    add_rect(sld, 10.2, 2.42, 2.88, 1.95, fill=NAVY)
    add_text(sld, '60% 효과 기준\n핵심 추정치', 10.32, 2.47, 2.64, 0.45,
             font_size=9, bold=True, color=RGBColor(0xA0,0xC4,0xFF), align=PP_ALIGN.CENTER)
    add_text(sld, '−21.3%p', 10.32, 2.93, 2.64, 0.55,
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sld, '보호자침해\n59.6%→38.3%', 10.32, 3.50, 2.64, 0.42,
             font_size=9, color=RGBColor(0xA0,0xC4,0xFF), align=PP_ALIGN.CENTER)

    # 연쇄 효과 추정
    add_rect(sld, 0.25, 3.90, 12.83, 0.68, fill=RGBColor(0xE8,0xF5,0xE9))
    add_text(sld, '연쇄 효과 추정 (보수적):',
             0.37, 3.93, 3.5, 0.25, font_size=9.5, bold=True, color=TEAL)
    add_text(sld,
             '악성민원패널티 60% 효과 → 보호자침해 −21.3%p → 자살사고 16.9% → 약 15.1%로 감소 가능 (보수적 추정)',
             3.9, 3.94, 9.0, 0.25, font_size=9.5, color=TEXT_DARK)
    add_text(sld,
             '수업방해분리 60% 효과 → 학생침해 −19.3%p → 이직고려율 46.3% → 약 40.2%로 감소 가능 (보수적 추정)',
             3.9, 4.22, 9.0, 0.25, font_size=9.5, color=TEXT_DARK)

    # 주의사항 박스
    add_rect(sld, 0.25, 4.72, 12.83, 1.12, fill=LIGHT_RED)
    add_text(sld, '시뮬레이션 한계 (필수 명시)',
             0.37, 4.76, 5.0, 0.28, font_size=10, bold=True, color=RED)
    caveats = [
        '① PAF의 OR은 D2 수요 기반 (침해 경험 후 측정) → 상한 추정치로 해석 권장 (인과 추론의 불확실성 존재)',
        '② Pe는 D2≥4 응답률로 추정 (전 교사 97-99% 수준 → 거의 포화 상태)',
        '③ 효과율(30/60/90%)은 정책 정착 단계 가정 — 실제 효과는 시행 환경에 따라 달라질 수 있음',
        '④ 단일 제도 효과만 고려 — 복합 정책 시 중복 효과 조정 필요',
    ]
    for i, cv in enumerate(caveats):
        add_text(sld, cv, 0.37, 5.06 + i*0.235, 12.5, 0.22,
                 font_size=9, color=TEXT_DARK)

    # 이미지
    img_path = BASE + 'output/final_story.png'
    if os.path.exists(img_path):
        add_image(sld, img_path, 0.25, 5.95, 12.83, 1.28)

    add_text(sld, '출처: D1(2024) | 방법: Rockhill et al. (1998) PAF | 분석: analysis_final_story.py',
             0.25, 7.28, 12.83, 0.20, font_size=7.5, color=TEXT_LIGHT, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — 결론
# ═══════════════════════════════════════════════════════════
def make_slide10(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, '10', '결론 및 정책 제언: 교권침해는 제도 설계의 실패다')

    # 핵심 결론
    add_rect(sld, 0.25, 0.72, 12.83, 0.55, fill=NAVY)
    add_text(sld,
             '"교권침해는 개인 문제도, 교육 부족도 아닌 학교 구조와 제도 설계의 실패다 — 맞춤 입법으로만 해결 가능하다"',
             0.37, 0.78, 12.6, 0.45,
             font_size=12, bold=True, color=WHITE_WARM, align=PP_ALIGN.CENTER)

    # 5대 발견 요약 테이블
    add_text(sld, '5대 핵심 발견 요약', 0.25, 1.35, 10.0, 0.28, font_size=11, bold=True, color=NAVY)
    findings = [
        ['#', '발견', '방법', '함의'],
        ['1', '인권교육 +20.8%p에도 만족도 3.75/5, 교사 47%가 제도 긴장 인식', 'D3 B2_6·B5_7', '교육이 아닌 제도가 답'],
        ['2', '초등 OR=1.79***, 학교규모 OR=1.29*** — 역인과 없는 구조적 원인', 'Track A', '구조적 취약 집단 확인'],
        ['3', '침해지수→이직 OR=2.17***, 자살 OR=2.57*** (직접효과, 구조통제)', 'Track A', '침해의 직접적 심각성'],
        ['4', '3집단: 완전히 다른 침해 구조 — Track A 패턴 독립 재현', 'K-means', '맞춤 입법 필요 확인'],
        ['5', 'C2 초등 관리자압박: D2 4개 모두 n.s. — 측정 도구 부재 확인', 'K-means+B', '현행 논의의 사각지대'],
    ]
    add_table(sld, findings, [0.4, 5.5, 1.8, 3.5],
              0.25, 1.65, 11.2, row_h=Inches(0.275),
              header_fill=NAVY, font_size=9)

    # 4대 정책 제언
    add_text(sld, '4대 정책 제언', 0.25, 3.72, 10.0, 0.28, font_size=11, bold=True, color=RED)
    policies = [
        ('즉시', '악성민원 법적패널티 입법', 'C1 중학교 최우선 — Track B Δ=+0.088*** / OR=2.380***', RED),
        ('단기', '수업방해 학생분리 조항 신설', 'C0 고등+C1 중학교 — Track B Δ=+0.128*** / SHAP 3위', TEAL),
        ('중기', '아동학대 신고 교원 보호 강화', 'C1 중학교 — OR=1.883***', ORANGE),
        ('신규', '관리자압박 내부신고 보호 시스템', 'C2 초등 전용 — 현행 법 논의에 없음 (신규 과제)', NAVY),
    ]
    for i, (timing, policy, rationale, color) in enumerate(policies):
        sy = 4.02 + i * 0.68
        add_rect(sld, 0.25, sy, 1.0, 0.55, fill=color)
        add_text(sld, timing, 0.25, sy+0.05, 1.0, 0.45,
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_rect(sld, 1.28, sy, 11.8, 0.55, fill=GRAY_BG if i % 2 == 0 else WHITE)
        add_text(sld, policy, 1.38, sy+0.04, 3.8, 0.26,
                 font_size=10.5, bold=True, color=color)
        add_text(sld, rationale, 1.38, sy+0.30, 11.5, 0.22,
                 font_size=9, color=TEXT_MED)

    # 최종 종합 메시지
    add_rect(sld, 0.25, 6.78, 12.83, 0.60, fill=NAVY)
    add_text(sld,
             '일률적 교권보호법이 아닌 학교급·침해유형에 따른 맞춤형 입법 로드맵 — '
             'C0(고등)·C1(중학교)·C2(초등) 각 집단의 구조와 수요가 다르다',
             0.37, 6.82, 12.5, 0.52,
             font_size=11, bold=True, color=WHITE_WARM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# SLIDE 11 — 참고문헌
# ═══════════════════════════════════════════════════════════
def make_reference(prs):
    sld = prs.slides.add_slide(BLANK)
    add_rect(sld, 0, 0, 13.33, 7.5, fill=WHITE_WARM)
    slide_header(sld, 'REF', '데이터 출처 및 참고문헌', accent=TEXT_MED)

    refs_left = [
        '【주요 데이터】',
        'D1. 윤정향. 교원 인권상황 실태조사, 2024 [데이터세트]. 한국노동사회연구소·국가인권위원회 / KOSSDA, 2025.',
        '    https://doi.org/10.22687/KOSSDA-A1-2024-0073-V1  (n=10,888)',
        '',
        'D3. 김병준. 초·중등 교원 인권교육 실태조사, 2021 [데이터세트]. 충남연구원·국가인권위원회 / KOSSDA, 2022.',
        '    https://doi.org/10.22687/KOSSDA-A1-2021-0019-V1.0  (n=9,553)',
        '',
        '한국교원단체총연합회. (2019). 2018년도 교권회복 및 교직상담 활동 실적 보고서. 한국교총 보도자료.',
        '국토교통부. (2024). 수도권 신도시 지구 현황. 도시정책과 공개자료.',
        '통계청. (2023). 국내인구이동통계 시군구별 순이동. KOSIS.',
        '교육부. (2023). 학원·교습소·개인과외 교습자 현황. 교육통계서비스 (EDSS).',
    ]
    refs_right = [
        '【분석 방법론】',
        'Igolkina & Meshcheryakov (2020). semopy: Python Package for SEM. Structural Equation Modeling.',
        'Lundberg & Lee (2017). A Unified Approach to Interpreting Model Predictions. NeurIPS.',
        'Brunner & Munzel (2000). The Nonparametric Behrens-Fisher Problem. Biometrical Journal.',
        'Preacher & Hayes (2008). Asymptotic and resampling strategies. Behavior Research Methods.',
        'Rockhill et al. (1998). Use and Misuse of Population Attributable Fractions. AJPH.',
        'MacQueen (1967). Methods for classification and analysis. 5th Berkeley Symposium.',
        '',
        '【분석 도구】',
        'Python 3.x: pyreadstat, pandas, numpy, scikit-learn, statsmodels, semopy, matplotlib',
        '분석 파일: analysis_two_track.py · analysis_sem.py · analysis_clustering.py · analysis_factor_impact.py',
    ]

    for i, ref in enumerate(refs_left):
        col = NAVY if ref.startswith('【') else TEXT_DARK
        bd = ref.startswith('【')
        add_text(sld, ref, 0.3, 0.82 + i * 0.50, 6.3, 0.47,
                 font_size=8.5 if not bd else 10, bold=bd, color=col)

    for i, ref in enumerate(refs_right):
        col = NAVY if ref.startswith('【') else TEXT_DARK
        bd = ref.startswith('【')
        add_text(sld, ref, 6.8, 0.82 + i * 0.50, 6.3, 0.47,
                 font_size=8.5 if not bd else 10, bold=bd, color=col)

    add_rect(sld, 6.65, 0.72, 0.03, 6.45, fill=GRAY_ROW)


# ═══════════════════════════════════════════════════════════
# 실행
# ═══════════════════════════════════════════════════════════
print("PPT 생성 시작...")
make_cover(prs)
print("  [커버] 완료")
make_slide01(prs)
print("  [01] 문제 제기 완료")
make_slide02(prs)
print("  [02] 데이터·방법 완료")
make_slide03(prs)
print("  [03] 인권교육 역설 완료")
make_slide04(prs)
print("  [04] Track A 구조→침해 완료")
make_slide05(prs)
print("  [05] Two-Track 수렴 완료")
make_slide06(prs)
print("  [06] K-means 클러스터링 완료")
make_slide07(prs)
print("  [07] C2 사각지대 완료")
make_slide08(prs)
print("  [08] 집단별 처방 완료")
make_slide09(prs)
print("  [09] PAF 시뮬레이션 완료")
make_slide10(prs)
print("  [10] 결론 완료")
make_reference(prs)
print("  [REF] 참고문헌 완료")

prs.save(OUT)
print(f"\n✓ PPT 저장 완료: {OUT}")
print(f"  총 {len(prs.slides)}슬라이드 (표지1 + 본문10 + 참고문헌1)")
