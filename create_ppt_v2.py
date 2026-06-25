"""
KOSSDA 2026 공모전 PPT v2 — 수상작 스타일
핵심: 차트가 슬라이드의 70%, 텍스트는 보조
대상 스타일: "0X | 제목" 헤더 + 차트 중심 + 데이터 수치 강조
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
from pptx.oxml.ns import qn
import os

BASE = '/Users/baiohelseu/Desktop/Project/kossda/'
IMGS = BASE + 'output/slides/'
OUT  = BASE + 'output/KOSSDA_2026_교권침해_v2.pptx'

# ── 색상 ─────────────────────────────────────────────────
NAVY    = RGBColor(0x0A, 0x24, 0x63)
NAVY_L  = RGBColor(0x1B, 0x4F, 0x8A)
RED     = RGBColor(0xC0, 0x39, 0x2B)
TEAL    = RGBColor(0x00, 0x7B, 0x7D)
ORANGE  = RGBColor(0xD3, 0x5A, 0x00)
GRAY    = RGBColor(0x8A, 0x8D, 0x9E)
LGRAY   = RGBColor(0xE8, 0xEC, 0xF0)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE= RGBColor(0xFA, 0xFA, 0xFB)
BLACK   = RGBColor(0x12, 0x14, 0x20)
REDBG   = RGBColor(0xFD, 0xED, 0xEB)
BLUEBG  = RGBColor(0xE3, 0xF0, 0xFB)
TEALBG  = RGBColor(0xE0, 0xF4, 0xF4)

KR = "Apple SD Gothic Neo"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── 헬퍼 ─────────────────────────────────────────────────
def rect(sld, x, y, w, h, fill=None, line=None, lw=Pt(0)):
    sh = sld.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill: sh.fill.solid(); sh.fill.fore_color.rgb = fill
    else:     sh.fill.background()
    if line:  sh.line.color.rgb = line; sh.line.width = lw
    else:     sh.line.fill.background()
    return sh

def txt(sld, text, x, y, w, h, size=11, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = sld.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.name = KR; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic; r.font.color.rgb = color
    return tb

def img(sld, path, x, y, w, h=None):
    if os.path.exists(path):
        if h: sld.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
        else: sld.shapes.add_picture(path, Inches(x), Inches(y), width=Inches(w))

def tbl(sld, data, col_w, x, y, total_w, row_h=0.28,
        hdr_fill=NAVY, hdr_txt=WHITE, alt=LGRAY, fs=8.5):
    rows, cols = len(data), len(data[0])
    t = sld.shapes.add_table(rows, cols,
        Inches(x), Inches(y), Inches(total_w), Inches(row_h*rows)).table
    for ci, cw in enumerate(col_w):
        t.columns[ci].width = Inches(cw)
    for ri, row in enumerate(data):
        for ci, val in enumerate(row):
            cell = t.cell(ri, ci)
            cell.text = ''
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            r = p.add_run(); r.text = val
            r.font.name = KR; r.font.size = Pt(fs)
            r.font.bold = (ri == 0); r.font.color.rgb = hdr_txt if ri == 0 else BLACK
            tc  = cell._tc
            tcP = tc.get_or_add_tcPr()
            sf  = etree.SubElement(tcP, qn('a:solidFill'))
            sc  = etree.SubElement(sf,  qn('a:srgbClr'))
            if ri == 0:     sc.set('val', str(hdr_fill))
            elif ri%2==0:   sc.set('val', str(alt))
            else:           sc.set('val', 'FFFFFF')


# ── 수상작 스타일 헤더 ───────────────────────────────────
def header(sld, num, title, bg=None, fg=None):
    """대상 수상작 스타일: 큰 섹션번호 + 제목 + 하단 구분선"""
    bg = bg or OFFWHITE
    fg = fg or NAVY
    rect(sld, 0, 0, 13.33, 7.5, fill=bg)        # 슬라이드 배경
    txt(sld, num,   0.25, 0.08, 0.70, 0.68, size=44, bold=True, color=fg)
    # 구분선
    rect(sld, 1.05, 0.22, 0.04, 0.42, fill=fg)  # 수직 구분선
    txt(sld, title, 1.18, 0.10, 11.9, 0.60, size=18, bold=True, color=fg)
    rect(sld, 0.25, 0.78, 12.83, 0.025, fill=LGRAY)  # 가로 구분선


def big_stat(sld, val, label, sub, x, y, w=2.8, h=1.1,
             bg=BLUEBG, val_c=NAVY, lbl_c=BLACK, sub_c=GRAY):
    """수상작 스타일 대형 통계 박스"""
    rect(sld, x, y, w, h, fill=bg)
    txt(sld, val,   x+0.08, y+0.06, w-0.16, h*0.48, size=30, bold=True, color=val_c, align=PP_ALIGN.CENTER)
    txt(sld, label, x+0.08, y+h*0.50, w-0.16, h*0.28, size=10, bold=True, color=lbl_c, align=PP_ALIGN.CENTER)
    txt(sld, sub,   x+0.08, y+h*0.78, w-0.16, h*0.20, size=8,  color=sub_c, align=PP_ALIGN.CENTER)


def note_box(sld, text, x, y, w, h, bg=BLUEBG, bc=NAVY_L, tc=BLACK, ts=9.5):
    """핵심 메시지 박스"""
    sh = rect(sld, x, y, w, h, fill=bg, line=bc, lw=Pt(1.5))
    txt(sld, text, x+0.12, y+0.10, w-0.22, h-0.18, size=ts, color=tc, wrap=True)


# ═══════════════════════════════════════════════════════════
# COVER
# ═══════════════════════════════════════════════════════════
def make_cover(prs):
    sld = prs.slides.add_slide(BLANK)
    rect(sld, 0, 0, 13.33, 7.5, fill=NAVY)
    rect(sld, 0, 0, 13.33, 0.15, fill=RED)
    rect(sld, 0, 7.35, 13.33, 0.15, fill=RED)

    # 배지
    rect(sld, 0.45, 0.35, 3.5, 0.38, fill=NAVY_L)
    txt(sld, 'KOSSDA 대학생 데이터 시각화 공모전 2026',
        0.5, 0.38, 3.4, 0.32, size=8.5, bold=True,
        color=RGBColor(0xA0,0xC4,0xFF), align=PP_ALIGN.CENTER)

    # 메인 제목
    txt(sld,
        '"인권교육 참여율이 20%p 증가했는데도\n교권침해는 왜 줄지 않는가?"',
        0.45, 0.90, 12.43, 2.20,
        size=36, bold=True, color=WHITE)

    # 구분선
    rect(sld, 0.45, 3.25, 7.5, 0.05, fill=RED)
    txt(sld, '— 제도 공백의 구조, 취약 집단, 그리고 맞춤 처방',
        0.45, 3.40, 12.43, 0.50,
        size=17, color=RGBColor(0xB0,0xD0,0xFF))

    # 4개 핵심 통계
    stats = [('85.6%','교권침해 경험률',BLUEBG, NAVY),
             ('OR=2.57','침해→자살사고',REDBG,  RED),
             ('47.3%','교사 제도긴장 인식',TEALBG, TEAL),
             ('−21.3%p','악성민원 입법 시\n침해율 감소',RGBColor(0xFF,0xF8,0xE1),ORANGE)]
    for i, (v, l, bg, vc) in enumerate(stats):
        sx = 0.45 + i*3.22
        rect(sld, sx, 4.10, 3.08, 0.95, fill=bg)
        txt(sld, v, sx+0.1, 4.12, 2.88, 0.52, size=26, bold=True, color=vc, align=PP_ALIGN.CENTER)
        txt(sld, l, sx+0.1, 4.68, 2.88, 0.34, size=9,  color=BLACK, align=PP_ALIGN.CENTER)

    # 데이터/방법
    txt(sld, '주요 데이터: 교원 인권상황 실태조사 2024 (n=10,888) · 초·중등 교원 인권교육 실태조사 2021 (n=9,553)',
        0.45, 5.22, 12.43, 0.30, size=8.5, color=RGBColor(0x80,0xA0,0xD0))
    txt(sld, '분석 기법: K-means 클러스터링 · SEM · SHAP · Bootstrap 매개효과 · PAF 시뮬레이션 · Two-Track 로지스틱 회귀',
        0.45, 5.52, 12.43, 0.30, size=8.5, color=RGBColor(0x80,0xA0,0xD0))
    txt(sld, 'Two-Track 전략: Track A (역인과 없는 구조 변수 OR) + Track B (당사자 정책 수요 격차)',
        0.45, 5.82, 12.43, 0.30, size=8.5, color=RGBColor(0x80,0xA0,0xD0))


# ═══════════════════════════════════════════════════════════
# SLIDE 01 — 문제 제기
# ═══════════════════════════════════════════════════════════
def make_s01(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '01', '문제 제기: 교사가 병들고 있다')

    # 핵심 수치 3개 (대형)
    big_stat(sld, '85.6%', '교권침해 경험률', '1종 이상 (D1 2024, n=10,888)',
             0.25, 0.90, 3.0, 1.08, bg=BLUEBG, val_c=NAVY)
    big_stat(sld, '46.3%', '이직 고려율', '침해경험자 OR=2.17***',
             3.42, 0.90, 3.0, 1.08, bg=RGBColor(0xFF,0xF8,0xE1), val_c=ORANGE)
    big_stat(sld, '16.9%', '자살사고율', '침해경험자 OR=2.57***',
             6.59, 0.90, 3.0, 1.08, bg=REDBG, val_c=RED)
    big_stat(sld, '59.6%', '보호자 민원 경험', '침해 유형 중 최다',
             9.76, 0.90, 3.32, 1.08, bg=TEALBG, val_c=TEAL)

    # 차트 (전체 넓이)
    img(sld, IMGS+'chart01_problem.png', 0.25, 2.08, 12.83, 4.62)

    # 출처
    txt(sld, '출처: D1 교원 인권상황 실태조사(2024)  |  *** p<.001 구조 변수 통제 후',
        0.25, 6.78, 12.83, 0.25, size=8, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 02 — 데이터와 분석 방법
# ═══════════════════════════════════════════════════════════
def make_s02(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '02', '데이터 선택 및 분석 방법: Two-Track 전략으로 역인과를 해소한다')

    # 왼쪽: 데이터 테이블
    txt(sld, '데이터 소개', 0.25, 0.90, 6.3, 0.28, size=11, bold=True, color=NAVY)
    tbl(sld, [
        ['데이터명', 'n', '활용'],
        ['교원 인권상황 실태조사(2024)', '10,888', '주 분석'],
        ['초·중등 교원 인권교육 실태조사(2021)', '9,553', '시계열·효과성'],
        ['한국교총 교권상담 실적(2018)', '501건', '외부 보조'],
        ['신도시 현황+순이동인구', '32개 단지', '지역 구조'],
        ['학원·교습소 현황(시군구)', '시군구', '교육열 프록시'],
    ], [4.5, 1.2, 2.3], 0.25, 1.20, 8.0, row_h=0.295, fs=9)

    # 왼쪽: 분석 기법 테이블
    txt(sld, '분석 기법', 0.25, 3.22, 6.3, 0.28, size=11, bold=True, color=NAVY)
    tbl(sld, [
        ['기법', '목적'],
        ['Two-Track 로지스틱 회귀', '역인과 해소 핵심 (Track A+B)'],
        ['K-means (k=3, Sil=0.404)', '취약 집단 3유형 분류'],
        ['SEM (semopy)', '구조적 인과 경로 분석'],
        ['SHAP (RandomForest n=300)', '제도 변수 중요도 해석'],
        ['Brunner-Munzel', '비모수 집단 비교'],
        ['Bootstrap 매개효과 n=2,000', '간접효과 95%CI 추정'],
        ['PAF 시뮬레이션', '정책 효과 추정 (역학 표준)'],
    ], [3.8, 5.0], 0.25, 3.52, 8.8, row_h=0.275, fs=9)

    # 오른쪽: Two-Track 도식 (텍스트 기반)
    rect(sld, 9.1, 0.90, 4.0, 5.90, fill=RGBColor(0xF4,0xF6,0xF8))
    txt(sld, 'Two-Track 설계', 9.22, 0.96, 3.76, 0.30, size=12, bold=True, color=NAVY)
    rect(sld, 9.22, 1.30, 3.76, 0.03, fill=LGRAY)

    # Track A
    rect(sld, 9.22, 1.40, 3.76, 1.55, fill=BLUEBG)
    txt(sld, 'Track A', 9.34, 1.44, 3.52, 0.28, size=11, bold=True, color=NAVY)
    txt(sld,
        '구조 변수\n(학교급·규모·고용형태)\n\n         ↓  역인과 없음 ✓\n\n침해 유형별 로지스틱 OR',
        9.34, 1.72, 3.52, 1.18, size=10, color=NAVY)

    rect(sld, 9.22, 3.05, 3.76, 0.03, fill=NAVY_L)

    # Track B
    rect(sld, 9.22, 3.15, 3.76, 1.55, fill=REDBG)
    txt(sld, 'Track B', 9.34, 3.19, 3.52, 0.28, size=11, bold=True, color=RED)
    txt(sld,
        'D2 정책 수요\n(침해집단 vs 비침해집단)\n\n  역인과 → 당사자 수요로 전환 ✓\n\nΔ 격차 = 우선 입법 타깃',
        9.34, 3.47, 3.52, 1.18, size=10, color=RED)

    # 수렴
    rect(sld, 9.22, 4.80, 3.76, 0.60, fill=NAVY)
    txt(sld, '수렴 타당도\nA×B가 같은 제도를 지목 → 확정',
        9.34, 4.87, 3.52, 0.50, size=10, bold=True, color=WHITE)

    # D2 내생성
    rect(sld, 0.25, 6.35, 8.8, 0.55, fill=REDBG, line=RED, lw=Pt(1))
    txt(sld,
        'D2 내생성 처리: D2는 침해 경험 후 측정 → Track A에서 완전 제외 / Track B에서 "수요 근거"로만 사용\n'
        'PAF OR = 수요 기반 상한 추정치 (인과 주장 한계 명시)',
        0.37, 6.40, 8.6, 0.48, size=9, color=RGBColor(0x80,0x10,0x10))

    txt(sld, '출처: KOSSDA 공개 데이터 5종', 9.1, 6.82, 4.0, 0.20, size=8, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 03 — HOOK: 인권교육 역설
# ═══════════════════════════════════════════════════════════
def make_s03(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '03', '인권교육 역설 [HOOK]: 교육량이 늘었는데 왜 침해는 줄지 않는가?', bg=TEALBG, fg=TEAL)

    # 3개 핵심 수치
    big_stat(sld, '+20.8%p', '인권교육 참여율\n2021→2024 증가',  '67.4% → 88.2%',
             0.25, 0.90, 3.92, 1.05, bg=TEALBG, val_c=TEAL)
    big_stat(sld, '3.75/5', '교육 효과 만족도\n(D3 B2_6)',       '보통 이하 39.1%',
             4.33, 0.90, 4.35, 1.05, bg=REDBG, val_c=RED)
    big_stat(sld, '47.3%', '"학생인권↑→교권↓"\n동의 비율',        '제도적 긴장 인식 지속',
             8.84, 0.90, 4.24, 1.05, bg=RGBColor(0xFF,0xF8,0xE1), val_c=ORANGE)

    # 차트 메인
    img(sld, IMGS+'chart02_education_paradox.png', 0.25, 2.05, 12.83, 4.50)

    # 결론 박스
    note_box(sld,
             '결론: 교육 참여율이 20%p 올랐는데도 효과 만족도 3.75/5, 교사 47.3%가 제도 긴장 인식 → 인식 교육이 아닌 제도 설계의 문제',
             0.25, 6.62, 12.83, 0.52, bg=TEALBG, bc=TEAL, tc=RGBColor(0x00,0x40,0x40), ts=10.5)

    txt(sld, '출처: D3 초·중등 교원 인권교육 실태조사(2021) B1·B2_6·B5_7  /  D1(2024) 참여율',
        0.25, 7.22, 12.83, 0.22, size=7.5, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 04 — Track A: 구조 → 침해 (Forest Plot)
# ═══════════════════════════════════════════════════════════
def make_s04(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '04', '구조가 침해를 만든다 — Track A (역인과 없는 인과 추론)')

    # OR 핵심 수치 (오른쪽 상단)
    big_stat(sld, 'OR=1.79***', '초등→학생침해', '구조 변수만으로 확인',
             0.25, 0.90, 3.0, 0.95, bg=BLUEBG, val_c=NAVY)
    big_stat(sld, 'OR=1.29***', '학교규모↑→침해', '1단계 상승 시',
             3.42, 0.90, 3.0, 0.95, bg=TEALBG, val_c=TEAL)
    big_stat(sld, 'OR=2.57***', '침해→자살사고', '구조 통제 후 직접효과',
             6.59, 0.90, 3.0, 0.95, bg=REDBG, val_c=RED)
    big_stat(sld, 'OR=2.17***', '침해→이직고려', '구조 통제 후 직접효과',
             9.76, 0.90, 3.32, 0.95, bg=RGBColor(0xFF,0xF8,0xE1), val_c=ORANGE)

    # 히트맵 + 학교규모 차트
    img(sld, IMGS+'chart04_structure.png', 0.25, 1.95, 12.83, 3.60)

    # Forest Plot
    txt(sld, '구조 변수 로지스틱 OR (Forest Plot)', 0.25, 5.62, 10.0, 0.28,
        size=10.5, bold=True, color=NAVY)
    img(sld, IMGS+'chart03_track_a_forest.png', 0.25, 5.92, 12.83, 1.35)

    txt(sld, '출처: D1(2024) | 기준집단: 유치원·특수·기타  *** p<.001  ** p<.01  * p<.05  n.s. 비유의 | 역인과 없음: 구조 변수는 침해 이전에 결정',
        0.25, 7.33, 12.83, 0.18, size=7.5, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 05 — Track B: 정책 수요 격차
# ═══════════════════════════════════════════════════════════
def make_s05(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '05', 'Track B: 당사자 정책 수요 격차 — 역인과를 강점으로 전환한다')

    # 핵심 수치
    big_stat(sld, 'Δ=+0.088***', '보호자침해집단\n악성민원패널티 수요격차', 'Mann-Whitney 유의',
             0.25, 0.90, 4.0, 0.95, bg=REDBG, val_c=RED)
    big_stat(sld, 'Δ=+0.128***', '학생침해집단\n수업방해분리 수요격차',    'Mann-Whitney 최대',
             4.42, 0.90, 4.0, 0.95, bg=TEALBG, val_c=TEAL)
    big_stat(sld, '95~99%', '전 교사 수요율\n(D2≥4 응답)', '천장효과 → 의미: 모두 필요하다',
             8.84, 0.90, 4.24, 0.95, bg=BLUEBG, val_c=NAVY)

    # 차트 (메인)
    img(sld, IMGS+'chart05_track_b_demand.png', 0.25, 1.95, 12.83, 3.92)

    # Track B 해석 원칙
    note_box(sld,
             'Track B 해석 원칙: D2는 침해 경험 이후 측정 → "제도 없어서 침해" 역인과 성립 X\n'
             '→ "침해 경험자가 절실히 요구" = 현장 기반 최우선 입법 타깃 (강점으로 전환)',
             0.25, 5.95, 8.5, 0.70, bg=REDBG, bc=RED, tc=RGBColor(0x70,0x10,0x10), ts=9.5)

    # 수렴 표
    tbl(sld, [
        ['침해 유형', 'Track A', 'Track B', '수렴 결론'],
        ['보호자침해(59.6%)', '초등 OR=1.47***', 'Δ악성민원=+0.088***', '악성민원패널티 1순위'],
        ['학생침해(78.5%)',   '규모 OR=1.29***', 'Δ수업분리=+0.128***', '수업방해분리 1순위'],
        ['관리자침해(23.6%)', 'D2 모두 비유의',  '격차 미약',            '신규 측정도구 선행'],
    ], [2.5, 2.2, 2.5, 3.2], 8.84, 5.95, 4.24, row_h=0.30, fs=8.5, hdr_fill=NAVY)

    txt(sld, '출처: D1(2024) | Track B는 인과 주장 없음 — "수요 근거"로만 해석',
        0.25, 7.33, 12.83, 0.18, size=7.5, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 06 — K-means 클러스터링
# ═══════════════════════════════════════════════════════════
def make_s06(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '06', '취약 집단 분류: 학교마다 다른 침해 구조 (K-means, k=3)')

    # 수치
    big_stat(sld, 'k=3', 'K-means\n최적 클러스터', 'Silhouette=0.404',
             0.25, 0.90, 2.8, 0.95, bg=BLUEBG, val_c=NAVY)
    big_stat(sld, '80.4%', 'C1 중학교형\n(최대 집단)', 'n=8,755 보호자침해 집중',
             3.22, 0.90, 2.8, 0.95, bg=TEALBG, val_c=TEAL)
    big_stat(sld, '7.8%', 'C2 초등형\n(사각지대)', 'n≈847 관리자침해 집중',
             6.19, 0.90, 2.8, 0.95, bg=REDBG, val_c=RED)
    big_stat(sld, '11.8%', 'C0 고등형\n(학생침해)', 'n=1,286',
             9.16, 0.90, 3.92, 0.95, bg=RGBColor(0xFF,0xF8,0xE1), val_c=ORANGE)

    # 레이더 차트 (메인)
    img(sld, IMGS+'chart06_kmeans_radar.png', 0.25, 1.95, 12.83, 4.05)

    # Track A 연계 검증
    note_box(sld,
             'Track A 수렴 검증: C1(중학교·대형) ↔ 학교규모 OR=1.29*** / C0(고등·학생) ↔ 고등 OR=0.66 보호방향\n'
             'C2(초등·소형) ↔ 초등 OR=1.47~1.79*** — K-means가 Track A 구조 패턴을 독립 재현 (수렴 타당도 확보)',
             0.25, 6.07, 12.83, 0.62, bg=BLUEBG, bc=NAVY, tc=NAVY, ts=10)

    txt(sld, '출처: D1(2024) | D2 제외, 순수 구조·침해 프로파일 클러스터링 | Silhouette Score=0.404',
        0.25, 7.33, 12.83, 0.18, size=7.5, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 07 — C2 사각지대
# ═══════════════════════════════════════════════════════════
def make_s07(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '07', '사각지대: 초등 관리자압박형은 현행 입법 논의에서 빠져 있다', bg=REDBG, fg=RED)

    # 핵심 수치
    big_stat(sld, '43.2%', 'C2 관리자침해율', '전체 평균 24.6% 대비 1.76배',
             0.25, 0.90, 3.92, 1.02, bg=REDBG, val_c=RED)
    big_stat(sld, '4개 모두\nn.s.', 'D2 현행 제도 OR', 'n≈847 검정력 제한',
             4.33, 0.90, 4.35, 1.02, bg=RGBColor(0xFF,0xF8,0xE1), val_c=ORANGE)
    big_stat(sld, '입법\n사각지대', '관리자압박 전용\n제도 없음', '현행 D2에 없는 유형',
             8.84, 0.90, 4.24, 1.02, bg=LGRAY, val_c=NAVY)

    # 차트 (메인)
    img(sld, IMGS+'chart07_blind_spot.png', 0.25, 2.00, 12.83, 3.92)

    # 2중 의미 박스
    rect(sld, 0.25, 6.00, 12.83, 1.10, fill=REDBG, line=RED, lw=Pt(1.5))
    txt(sld, '발견의 2중 의미:', 0.37, 6.06, 4.0, 0.26, size=10.5, bold=True, color=RED)
    texts_7 = [
        '① [입법 사각지대] 현행 학부모·학생 중심 교권보호법 논의는 초등 소규모 학교의 관리자 압박을 보호하지 못함',
        '② [측정 공백] n≈847 소집단 → 검정력 부족으로 "효과 없음"이 아닌 "검출 불가" → 측정 도구 개발 선행 필요',
        '③ [처방] 관리자 갑질·부당 지시 대응 내부신고·보호 시스템 신규 입법 → C2 전용 제도 필요',
    ]
    for i, t in enumerate(texts_7):
        txt(sld, t, 0.37, 6.36 + i*0.23, 12.5, 0.21, size=9.5,
            color=RED if i == 2 else BLACK, bold=(i == 2))

    txt(sld, '출처: D1(2024) | K-means C2 집단 (초등·관리자압박형) | *** p<.001  n.s. 비유의',
        0.25, 7.33, 12.83, 0.18, size=7.5, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 08 — Track A×B 수렴 + 정책 처방
# ═══════════════════════════════════════════════════════════
def make_s08(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '08', '집단별 맞춤 처방: Track A × B 수렴 — 일률적 교권보호법의 한계')

    # 수렴 차트 (상단 2/3)
    img(sld, IMGS+'chart09_convergence.png', 0.25, 0.90, 12.83, 4.10)

    # 처방 카드 3개 (하단 1/3)
    prescriptions = [
        ('C0\n고등형', '수업방해\n분리시스템', 'Track B Δ=+0.128***', NAVY),
        ('C1\n중학교형', '악성민원\n법적패널티', 'Track B Δ=+0.088***', TEAL),
        ('C2\n초등형', '관리자압박\n내부신고시스템', 'D2 n.s. → 신규입법', RED),
    ]
    for i, (cl, policy, rationale, color) in enumerate(prescriptions):
        sx = 0.25 + i * 4.36
        rect(sld, sx, 5.12, 4.15, 0.36, fill=color)
        txt(sld, cl, sx+0.1, 5.14, 1.0, 0.32, size=10.5, bold=True, color=WHITE)
        txt(sld, '→ ' + policy, sx+1.15, 5.14, 2.9, 0.32, size=11, bold=True, color=WHITE)
        rect(sld, sx, 5.48, 4.15, 0.55,
             fill=REDBG if color==RED else (TEALBG if color==TEAL else BLUEBG))
        txt(sld, rationale, sx+0.12, 5.52, 3.9, 0.48, size=10, color=color, bold=True)

    # 정책 로드맵 테이블
    tbl(sld, [
        ['단계', '대상', '정책', '시급도'],
        ['즉시', 'C1 중학교(80.4%)', '악성민원 법적패널티 입법', '★★★★★'],
        ['단기', 'C0·C1(학생·보호자)', '수업방해 학생분리 조항', '★★★★'],
        ['중기', 'C1 중학교', '아동학대 신고 교원보호 강화', '★★★'],
        ['신규', 'C2 초등(사각지대)', '관리자압박 내부신고 보호 시스템', '★★★(新)'],
    ], [0.8, 2.5, 5.5, 1.5], 0.25, 6.12, 10.3, row_h=0.27, fs=9, hdr_fill=NAVY)

    txt(sld, '출처: D1(2024) | analysis_two_track.py / analysis_clustering.py',
        0.25, 7.33, 12.83, 0.18, size=7.5, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 09 — PAF 시뮬레이션
# ═══════════════════════════════════════════════════════════
def make_s09(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '09', '정책 시뮬레이션: 제도 도입 시 침해율 감소 추정 (PAF)')

    # 핵심 수치
    big_stat(sld, '−21.3%p', '악성민원패널티 60%효과\n보호자침해 59.6%→38.3%', '수요 기반 상한 추정치',
             0.25, 0.90, 4.0, 1.02, bg=REDBG, val_c=RED)
    big_stat(sld, '−19.3%p', '수업방해분리 60%효과\n학생침해 78.5%→59.2%', '수요 기반 상한 추정치',
             4.42, 0.90, 4.0, 1.02, bg=TEALBG, val_c=TEAL)
    big_stat(sld, 'PAF', 'Population\nAttributable Fraction', '역학 표준 시뮬레이션 방법론',
             8.84, 0.90, 4.24, 1.02, bg=BLUEBG, val_c=NAVY)

    # 차트 (메인)
    img(sld, IMGS+'chart08_paf.png', 0.25, 2.00, 12.83, 3.92)

    # PAF 공식 + 한계
    rect(sld, 0.25, 6.00, 7.3, 0.55, fill=RGBColor(0xFF,0xF8,0xE1), line=ORANGE, lw=Pt(1))
    txt(sld, 'PAF = Pe × (OR−1) / [1 + Pe × (OR−1)]',
        0.37, 6.04, 7.1, 0.48, size=11, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

    rect(sld, 7.72, 6.00, 5.36, 0.55, fill=REDBG, line=RED, lw=Pt(1))
    txt(sld, '한계: OR은 D2 수요 기반 → 상한 추정치 | Pe = D2≥4 비율(97~99%) | 효과율 가정 포함',
        7.84, 6.04, 5.1, 0.48, size=9, color=RED)

    txt(sld, '출처: D1(2024) | 방법: Rockhill et al.(1998) PAF | *** p<.001 (D2 로지스틱)',
        0.25, 7.33, 12.83, 0.18, size=7.5, color=GRAY, italic=True)


# ═══════════════════════════════════════════════════════════
# SLIDE 10 — 결론
# ═══════════════════════════════════════════════════════════
def make_s10(prs):
    sld = prs.slides.add_slide(BLANK)
    header(sld, '10', '결론 및 정책 제언: 교권침해는 제도 설계의 실패다')

    # 핵심 결론 박스 (전체 폭)
    rect(sld, 0.25, 0.90, 12.83, 0.52, fill=NAVY)
    txt(sld,
        '"교권침해는 개인 문제도 교육 부족도 아닌, 학교 구조와 제도 설계의 실패다 — 맞춤 입법으로만 해결 가능하다"',
        0.37, 0.96, 12.6, 0.42, size=12.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 5대 발견 테이블
    txt(sld, '5대 핵심 발견', 0.25, 1.52, 8.0, 0.28, size=11, bold=True, color=NAVY)
    tbl(sld, [
        ['#', '발견', '방법', '함의'],
        ['1', '인권교육 +20.8%p에도 만족도 3.75/5, 47.3%가 제도 긴장 인식', 'D3 B2_6·B5_7', '교육이 아닌 제도'],
        ['2', '초등 OR=1.79***, 학교규모 OR=1.29*** — 역인과 없는 구조 원인', 'Track A', '구조 취약 집단 확인'],
        ['3', '침해 → 이직 OR=2.17***, 자살 OR=2.57*** (구조 통제 후 직접효과)', 'Track A', '침해의 직접 심각성'],
        ['4', '3집단: C0 고등학생형 / C1 중학교보호자형 / C2 초등관리자형', 'K-means', '맞춤 입법 필요 확인'],
        ['5', 'C2 초등형 D2 4개 n.s. → 현행 입법 논의의 측정·입법 사각지대', 'K-means+B', '신규 제도 시급'],
    ], [0.35, 5.8, 1.6, 3.2], 0.25, 1.82, 11.2, row_h=0.265, fs=8.5, hdr_fill=NAVY)

    # 4대 정책 제언
    txt(sld, '4대 정책 제언', 0.25, 3.70, 8.0, 0.28, size=11, bold=True, color=RED)
    policies = [
        ('즉시', '악성민원 법적패널티 입법', 'C1 중학교(80.4%) 최우선 — Track B Δ=+0.088*** / OR=2.380***', NAVY),
        ('단기', '수업방해 학생분리 조항', 'C0·C1 — Track B Δ=+0.128*** / SHAP 3위', TEAL),
        ('중기', '아동학대 신고 교원보호 강화', 'C1 — OR=1.883***', NAVY_L),
        ('신규', '관리자압박 내부신고 보호 시스템', 'C2 초등 전용 — D2 4개 n.s. → 측정·입법 사각지대 해소', RED),
    ]
    for i, (timing, policy, rationale, color) in enumerate(policies):
        sy = 4.00 + i * 0.68
        rect(sld, 0.25, sy, 0.80, 0.58, fill=color)
        txt(sld, timing, 0.25, sy+0.08, 0.80, 0.42, size=10.5, bold=True,
            color=WHITE, align=PP_ALIGN.CENTER)
        rect(sld, 1.10, sy, 12.0, 0.58,
             fill=LGRAY if i % 2 == 0 else RGBColor(0xFA,0xFA,0xFB))
        txt(sld, policy,    1.22, sy+0.03, 4.0, 0.26, size=11, bold=True, color=color)
        txt(sld, rationale, 1.22, sy+0.31, 11.7, 0.24, size=9, color=GRAY)

    # 최종
    rect(sld, 0.25, 6.78, 12.83, 0.57, fill=NAVY)
    txt(sld,
        '일률적 교권보호법이 아닌 학교급·침해유형별 맞춤 입법 로드맵 — C0(고등)·C1(중학교)·C2(초등) 각 집단의 구조와 수요가 다르다',
        0.37, 6.84, 12.6, 0.48, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# 실행
# ═══════════════════════════════════════════════════════════
print("PPT v2 생성 중...")
make_cover(prs)
print("  [표지] ✓")
make_s01(prs); print("  [01] 문제제기 ✓")
make_s02(prs); print("  [02] 데이터·방법 ✓")
make_s03(prs); print("  [03] 인권교육 역설 ✓")
make_s04(prs); print("  [04] Track A ✓")
make_s05(prs); print("  [05] Track B ✓")
make_s06(prs); print("  [06] K-means ✓")
make_s07(prs); print("  [07] C2 사각지대 ✓")
make_s08(prs); print("  [08] 수렴·처방 ✓")
make_s09(prs); print("  [09] PAF ✓")
make_s10(prs); print("  [10] 결론 ✓")

prs.save(OUT)
print(f"\n✓ 저장: {OUT}")
print(f"  총 {len(prs.slides)}슬라이드")
